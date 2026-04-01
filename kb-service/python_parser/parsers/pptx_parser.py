"""
PPTX 解析器，基于 python-pptx
"""
import logging
from pptx import Presentation

from utils.common import split_into_chunks

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str, doc_id: str) -> dict:
    """
    解析 PPTX 文件，返回与 Go ParsedDocument 对应的 dict。
    策略：按幻灯片逐页提取文字，保留页码和幻灯片标题。
    """
    prs = Presentation(file_path)
    total_slides = len(prs.slides)

    slide_parts = []
    full_text_parts = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # 提取幻灯片标题（第一个形状中的标题占位符）
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.shape_type.name == "TITLE" or shape.shape_type.name == "CLOSING_TITLE_TEXT":
                    title_text = shape.text_frame.text.strip()
                    break

        # 提取正文
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        combined = "\n".join(texts)
        if not combined and not title_text:
            continue

        slide_parts.append({
            "slide_number": slide_num,
            "title": title_text,
            "text": combined,
        })
        if combined:
            full_text_parts.append(combined)

    # 全量文本分块
    all_text = "\n\n".join(full_text_parts)
    raw_chunks = split_into_chunks(all_text, doc_id)
    chunks = _attach_slide_info(raw_chunks, slide_parts)

    title = _extract_title(slide_parts)
    summary = _build_summary(full_text_parts)

    return {
        "doc_id": doc_id,
        "file_type": "pptx",
        "title": title,
        "text_chunks": chunks,
        "summary": summary,
        "total_pages": total_slides,
    }


def _attach_slide_info(chunks: list[dict], slide_parts: list[dict]) -> list[dict]:
    """为每个 chunk 补充页码和章节标题"""
    if not slide_parts:
        return chunks

    # 建立字符偏移到幻灯片信息的映射
    boundaries = []
    offset = 0
    for sp in slide_parts:
        boundaries.append((offset, sp["slide_number"], sp["title"]))
        offset += len(sp["text"]) + 2

    for chunk in chunks:
        start = chunk["metadata"]["start_char"]
        page_num, section = 1, ""
        for boundary, pnum, title in boundaries:
            if start >= boundary:
                page_num, section = pnum, title
        chunk["metadata"]["page_number"] = page_num
        chunk["metadata"]["section_title"] = section

    return chunks


def _extract_title(slide_parts: list[dict]) -> str:
    for sp in slide_parts:
        if sp["title"]:
            return sp["title"][:80]
    return "PPT 文档"


def _build_summary(full_text_parts: list[str]) -> str:
    combined = "\n".join(full_text_parts)
    return combined[:200] + "..." if len(combined) > 200 else combined
