#!/usr/bin/env python3
"""
PPT Page Renderer

Receives a JSON payload via stdin with:
{
    "page_index": int,           # 0-based page index
    "page_title": str,           # Title of this page
    "output_path": str,          # Where to write the .pptx file
    "py_code": str,              # Python code to execute
    "render_config": {
        "width_inches": float,   # Slide width (default 10)
        "height_inches": float,  # Slide height (default 7.5)
        "bg_color": str,         # Background color hex (default "FFFFFF")
        "font_name": str,        # Default font (default "Microsoft YaHei")
    }
}

Writes result to stdout as JSON:
{
    "success": bool,
    "pptx_path": str,            # Absolute path to generated .pptx
    "render_url": str,           # URL where file will be served (empty on failure)
    "error": str                 # Error message (empty on success)
}
"""

import json
import sys
import traceback


def run():
    raw = sys.stdin.read()
    if not raw.strip():
        print(json.dumps({"success": False, "pptx_path": "", "render_url": "", "error": "empty input"}))
        sys.exit(1)

    try:
        req = json.loads(raw)
    except json.JSONDecodeError as e:
        print(json.dumps({"success": False, "pptx_path": "", "render_url": "", "error": f"invalid json: {e}"}))
        sys.exit(1)

    page_index = req.get("page_index", 0)
    page_title = req.get("page_title", "")
    output_path = req.get("output_path", "")
    py_code = req.get("py_code", "")
    cfg = req.get("render_config", {})

    width = cfg.get("width_inches", 10)
    height = cfg.get("height_inches", 7.5)
    bg_color = cfg.get("bg_color", "FFFFFF")
    font_name = cfg.get("font_name", "Microsoft YaHei")

    if not output_path:
        print(json.dumps({"success": False, "pptx_path": "", "render_url": "", "error": "output_path is required"}))
        sys.exit(1)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import copy

        prs = Presentation()
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)

        bg = RGBColor.from_string(bg_color.lstrip("#"))
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

        def rgb(hex_color: str) -> RGBColor:
            h = hex_color.lstrip("#")
            return RGBColor.from_string(h)

        def emu(inches: float) -> Emu:
            return Inches(inches)

        def pt(size: float) -> Pt:
            return Pt(size)

        def set_paragraph(shape, text, left=0.5, top=0.5, width=9, height=1,
                          font_size=18, bold=False, color="000000",
                          align=PP_ALIGN.LEFT, font_name=font_name):
            from pptx.util import Inches, Pt
            shape.text_frame.clear()
            shape.text_frame.word_wrap = True
            p = shape.text_frame.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
            run.font.name = font_name

        def add_textbox(slide, text, left, top, width, height,
                        font_size=18, bold=False, color="000000",
                        align=PP_ALIGN.LEFT):
            from pptx.util import Inches, Pt
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                              Inches(width), Inches(height))
            set_paragraph(txBox, text, left, top, width, height,
                          font_size, bold, color, align)

        def add_rect(slide, left, top, width, height, fill="FFFFFF", line="CCCCCC"):
            from pptx.util import Inches
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
            if line and line != "none":
                shape.line.color.rgb = rgb(line)
                shape.line.width = Pt(0.5)
            else:
                shape.line.fill.background()
            return shape

        def add_oval(slide, left, top, width, height, fill="4472C4", line="none"):
            shape = slide.shapes.add_shape(
                9,  # MSO_SHAPE_TYPE.OVAL
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
            if line and line != "none":
                shape.line.color.rgb = rgb(line)
            else:
                shape.line.fill.background()
            return shape

        def add_image(slide, path, left, top, width, height):
            from pptx.util import Inches
            pic = slide.shapes.add_picture(path, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
            return pic

        def set_slide_title(slide, text, font_size=32, color="FFFFFF",
                           bg_color="1F4E79", height=1.2):
            from pptx.util import Inches, Pt
            title_box = slide.shapes.add_textbox(
                Inches(0), Inches(0), Inches(width), Inches(height)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = rgb(color)
            run.font.name = font_name

        # Build globals for user code
        user_globals = {
            "prs": prs,
            "slide": slide,
            "width": width,
            "height": height,
            "rgb": rgb,
            "emu": emu,
            "pt": pt,
            "add_textbox": add_textbox,
            "add_rect": add_rect,
            "add_oval": add_oval,
            "add_image": add_image,
            "set_slide_title": set_slide_title,
            "PP_ALIGN": PP_ALIGN,
            "font_name": font_name,
        }
        user_globals.update({
            "add_shape": slide.shapes.add_shape,
            "add_table": slide.shapes.add_table,
        })

        exec_globals = {"__builtins__": {}}
        exec_globals.update(user_globals)

        if py_code.strip():
            exec(py_code, exec_globals)

        prs.save(output_path)

        render_url = req.get("render_url_prefix", "") + output_path
        print(json.dumps({
            "success": True,
            "pptx_path": output_path,
            "render_url": render_url,
            "error": ""
        }))

    except Exception as e:
        tb = traceback.format_exc()
        print(json.dumps({
            "success": False,
            "pptx_path": "",
            "render_url": "",
            "error": f"{type(e).__name__}: {e}\n{tb}"
        }))
        sys.exit(1)


if __name__ == "__main__":
    run()
