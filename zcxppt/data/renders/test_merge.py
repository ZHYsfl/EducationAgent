import sys
sys.path.insert(0, 'e:/PPT agent/EducationAgent/zcxppt/internal/infra/renderer')
import json, os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(hex_color):
    return RGBColor.from_string(hex_color.lstrip('#'))

def add_rect(slide, left, top, width, height, fill='FFFFFF', line='CCCCCC'):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line and line != 'none':
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False, color='000000', align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = 'Microsoft YaHei'

def set_slide_title(slide, text, font_size=32, color='FFFFFF', bg_color='1F4E79', height=1.2):
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(height))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)
    run.font.name = 'Microsoft YaHei'

page_codes = {
    'page_1': 'add_rect(slide, 0, 0, 10, 1.2, fill="1F4E79")\nset_slide_title(slide, "Welcome to AI Basics", font_size=36, color="FFFFFF", bg_color="1F4E79")\nadd_textbox(slide, "Today we will learn about:", 0.5, 1.5, 9, 0.5, font_size=20, bold=True, color="333333")\nadd_textbox(slide, "1. What is Artificial Intelligence?", 1.0, 2.2, 7, 0.5, font_size=16, color="555555")\nadd_textbox(slide, "2. Machine Learning Basics", 1.0, 2.8, 7, 0.5, font_size=16, color="555555")',
    'page_2': 'add_rect(slide, 0, 0, 10, 1.2, fill="2E7D32")\nset_slide_title(slide, "Machine Learning Basics", font_size=36, color="FFFFFF", bg_color="2E7D32")\nadd_textbox(slide, "Machine learning is a subset of AI that enables systems to learn from data.", 0.5, 1.5, 9, 1.0, font_size=18, color="333333")\nadd_textbox(slide, "Key Types of Machine Learning:", 0.5, 2.6, 9, 0.5, font_size=20, bold=True, color="1F4E79")\nadd_textbox(slide, "Supervised Learning - Learning with labeled data", 1.0, 3.2, 8, 0.5, font_size=16, color="555555")\nadd_textbox(slide, "Unsupervised Learning - Finding patterns in unlabeled data", 1.0, 3.8, 8, 0.5, font_size=16, color="555555")\nadd_textbox(slide, "Reinforcement Learning - Learning through interaction", 1.0, 4.4, 8, 0.5, font_size=16, color="555555")',
    'page_3': 'add_rect(slide, 0, 0, 10, 1.2, fill="7B1FA2")\nset_slide_title(slide, "Summary and Questions", font_size=36, color="FFFFFF", bg_color="7B1FA2")\nadd_textbox(slide, "What we covered today:", 0.5, 1.5, 9, 0.5, font_size=20, bold=True, color="333333")\nadd_textbox(slide, "What is Artificial Intelligence", 1.0, 2.2, 8, 0.5, font_size=16, color="555555")\nadd_textbox(slide, "Machine Learning and its types", 1.0, 2.8, 8, 0.5, font_size=16, color="555555")\nadd_textbox(slide, "Do you have any questions?", 0.5, 4.0, 9, 0.8, font_size=24, bold=True, color="1F4E79", align=PP_ALIGN.CENTER)',
}

os.makedirs('e:/PPT agent/EducationAgent/zcxppt/data/renders', exist_ok=True)

merged_prs = Presentation()
merged_prs.slide_width = Inches(10)
merged_prs.slide_height = Inches(7.5)

ordered_ids = sorted(page_codes.keys())

for page_id in ordered_ids:
    code = page_codes[page_id]
    blank_layout = merged_prs.slide_layouts[6]
    slide = merged_prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb('FFFFFF')

    user_globals = {'add_rect': add_rect, 'add_textbox': add_textbox, 'set_slide_title': set_slide_title,
                   'slide': slide, 'prs': merged_prs, 'PP_ALIGN': PP_ALIGN, 'font_name': 'Microsoft YaHei'}
    exec_globals = {'__builtins__': {}}
    exec_globals.update(user_globals)
    if code and code.strip():
        exec(code, exec_globals)

output_path = 'e:/PPT agent/EducationAgent/zcxppt/data/renders/test_merged.pptx'
merged_prs.save(output_path)
print('SUCCESS: merged', len(ordered_ids), 'pages to', output_path)
print('File size:', os.path.getsize(output_path), 'bytes')
print('Slide count:', len(merged_prs.slides))
