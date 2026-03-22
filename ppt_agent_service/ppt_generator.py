from __future__ import annotations

import copy
import json
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Optional

import sys
import urllib.parse

# 让本服务可以直接 import PPTAgent 里的 python 包（`pptagent`）
_REPO_ROOT = Path(__file__).resolve().parents[1]
_PPTAGENT_ROOT = _REPO_ROOT / "PPTAgent"
if str(_PPTAGENT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PPTAGENT_ROOT))

from pptagent.document.document import Document
from pptagent.llms import AsyncLLM
from pptagent.presentation import Presentation
from pptagent.pptgen import PPTAgent
from pptagent.utils import Config, package_join, ppt_to_images

from ppt_agent_service.slide_py_code import extract_slide_html_from_py


def description_to_markdown(description: str, topic: str) -> str:
    """
    将 VoiceAgent 的 buildDetailedDescription 文本（包含 【字段】行）转成 markdown headings，
    以便 PPTAgent 的 Document.from_markdown 能按标题层级分块。
    """
    # 例如：【课程主题】xxx -> ## 课程主题\nxxx
    def repl(m: re.Match[str]) -> str:
        key = m.group(1).strip()
        val = m.group(2).strip()
        return f"## {key}\n{val}\n"

    # 情况 A：同一行，例如 【课程主题】xxx
    pattern_inline = r"【([^】]+)】\s*([^\n]+)"
    converted = re.sub(pattern_inline, repl, description)
    # 情况 B：换行，例如 【参考资料1：pdf】\n使用说明：...
    converted = re.sub(r"【([^】]+)】\s*\n", r"## \1\n", converted)
    # 兜底：没有任何字段标记时，直接作为正文
    if converted.strip() == description.strip():
        converted = description

    return f"# {topic}\n\n{converted.strip()}\n"


@dataclass
class GeneratedDeck:
    task_dir: Path
    pptx_path: Path
    slide_image_dir: Path
    slide_html: list[str]
    slide_count: int


class PPTGenerator:
    def __init__(
        self,
        template_name: str = "default",
        template_hidden_small_pic_ratio: float = 0.2,
    ):
        # 模型通过环境变量注入（与 PPTAgentServer 保持一致）
        model_name = os.getenv("PPTAGENT_MODEL")
        api_base = os.getenv("PPTAGENT_API_BASE")
        api_key = os.getenv("PPTAGENT_API_KEY")

        if not model_name:
            raise RuntimeError(
                "PPTAGENT_MODEL not set. Please set PPTAGENT_MODEL/PPTAGENT_API_BASE/PPTAGENT_API_KEY"
            )

        # 当前版本 PPTAgent 用同一个模型跑语言+视觉（足够用于最小 happy path）
        llm = AsyncLLM(model_name, api_base, api_key)

        self.language_model = llm
        self.vision_model = llm
        self.template_name = template_name
        self.template_hidden_small_pic_ratio = template_hidden_small_pic_ratio
        self.kb_service_url = os.getenv("KB_SERVICE_URL", "http://localhost:9200")
        self.search_service_url = os.getenv(
            "WEB_SEARCH_SERVICE_URL", "http://localhost:9400"
        )
        self.enable_web_search = os.getenv("PPTAGENT_ENABLE_WEB_SEARCH", "true").lower() == "true"
        self.kb_top_k = int(os.getenv("PPTAGENT_KB_TOPK", "8"))
        self.search_top_k = int(os.getenv("PPTAGENT_SEARCH_TOPK", "5"))

        self._template_presentation: Optional[Presentation] = None
        self._template_slide_induction: Optional[dict] = None

    def _load_template_once(self) -> None:
        if self._template_presentation is not None and self._template_slide_induction is not None:
            return

        templates_dir = Path(package_join("templates"))
        template_folder = templates_dir / self.template_name
        if not template_folder.exists():
            raise FileNotFoundError(f"PPT template not found: {self.template_name}")

        prs_config = Config(str(template_folder))
        prs = Presentation.from_file(
            str(template_folder / "source.pptx"), prs_config
        )
        slide_induction = json.loads(
            (template_folder / "slide_induction.json").read_text(encoding="utf-8")
        )

        self._template_presentation = prs
        self._template_slide_induction = slide_induction

    def _sanitize_filename(self, name: str) -> str:
        name = name.strip()
        # Only keep safe chars to avoid markdown / filesystem issues.
        name = re.sub(r"[^0-9a-zA-Z._-]+", "_", name)
        name = name.lstrip("._-")
        return name or "file"

    def _infer_extension(self, file_url: str, fallback: str = ".png") -> str:
        try:
            parsed = urllib.parse.urlparse(file_url)
            tail = Path(parsed.path).name
            ext = Path(tail).suffix
            if ext:
                return ext
        except Exception:
            pass
        return fallback

    async def _download_reference_image(
        self, *, file_url: str, dest_path: Path
    ) -> None:
        """
        Download/copy a reference image into local `dest_path`.
        """
        if dest_path.exists() and dest_path.stat().st_size > 0:
            return

        # Local path fallback
        if not file_url.startswith("http://") and not file_url.startswith("https://"):
            if file_url.startswith("file://"):
                src = Path(urllib.parse.urlparse(file_url).path)
            else:
                src = Path(file_url)
            if src.exists():
                dest_path.parent.mkdir(parents=True, exist_ok=True)
                dest_path.write_bytes(src.read_bytes())
                return

        import httpx

        dest_path.parent.mkdir(parents=True, exist_ok=True)
        async with httpx.AsyncClient(timeout=60.0) as client:
            r = await client.get(file_url)
            r.raise_for_status()
            dest_path.write_bytes(r.content)

    async def _call_api(self, method: str, url: str, json_body: dict | None = None) -> dict | None:
        import httpx

        try:
            async with httpx.AsyncClient(timeout=45.0) as client:
                if method == "GET":
                    resp = await client.get(url)
                else:
                    resp = await client.post(url, json=json_body or {})
                resp.raise_for_status()
                payload = resp.json()
                if isinstance(payload, dict) and "code" in payload:
                    if payload.get("code") != 200:
                        return None
                    return payload.get("data")
                return payload if isinstance(payload, dict) else None
        except Exception:
            return None

    async def _query_kb_chunks(self, *, user_id: str, query: str) -> list[dict]:
        data = await self._call_api(
            "POST",
            f"{self.kb_service_url}/api/v1/kb/query",
            {
                "user_id": user_id,
                "query": query,
                "top_k": self.kb_top_k,
                "score_threshold": 0.45,
            },
        )
        if not data or "chunks" not in data:
            return []
        chunks = data.get("chunks") or []
        return [c for c in chunks if isinstance(c, dict)]

    async def _search_web_summary(self, *, user_id: str, query: str) -> str:
        if not self.enable_web_search:
            return ""
        data = await self._call_api(
            "POST",
            f"{self.search_service_url}/api/v1/search/query",
            {
                "request_id": "",
                "user_id": user_id,
                "query": query,
                "max_results": self.search_top_k,
                "language": "zh",
                "search_type": "general",
            },
        )
        if not data:
            return ""
        summary = data.get("summary")
        return summary.strip() if isinstance(summary, str) else ""

    async def _parse_reference_file(self, *, file_url: str, file_type: str) -> dict | None:
        # Follows spec chapter 8 parse endpoint.
        return await self._call_api(
            "POST",
            f"{self.kb_service_url}/api/v1/kb/parse",
            {"file_url": file_url, "file_type": file_type, "doc_id": ""},
        )

    async def _build_external_context_markdown(
        self,
        *,
        user_id: str,
        topic: str,
        description: str,
        image_dir: Path,
        reference_files: list[dict] | None,
    ) -> str:
        """
        Gather multimodal materials from KB / web / provided references and convert into markdown.
        """
        blocks: list[str] = []
        query = f"{topic}\n{description[:1500]}"

        kb_chunks = await self._query_kb_chunks(user_id=user_id, query=query)
        if kb_chunks:
            blocks.append("## 本地知识库检索结果")
            for i, c in enumerate(kb_chunks, start=1):
                title = c.get("doc_title") or c.get("doc_id") or f"chunk_{i}"
                content = str(c.get("content") or "").strip()
                if not content:
                    continue
                blocks.append(f"### 知识片段{i}（来源：{title}）\n{content}")
                meta = c.get("metadata") or {}
                image_url = meta.get("image_url") if isinstance(meta, dict) else None
                if isinstance(image_url, str) and image_url.strip():
                    ext = self._infer_extension(image_url, fallback=".png")
                    local_name = self._sanitize_filename(f"kb_chunk_{i}{ext}")
                    try:
                        await self._download_reference_image(
                            file_url=image_url, dest_path=image_dir / local_name
                        )
                        blocks.append(f"![知识库配图{i}]({local_name})")
                    except Exception:
                        pass

        web_summary = await self._search_web_summary(user_id=user_id, query=topic)
        if web_summary:
            blocks.append("## 网络检索摘要")
            blocks.append(web_summary)

        if reference_files:
            blocks.append("## 参考资料解析结果")
            parse_jobs = []
            for rf in reference_files:
                ft = (rf.get("file_type") or "").lower()
                fu = (rf.get("file_url") or "").strip()
                if not fu or ft in ("image", "images"):
                    continue
                parse_jobs.append((rf, ft, fu))

            import asyncio

            async def _one_parse(rf: dict, ft: str, fu: str) -> tuple[dict, dict | None]:
                return rf, await self._parse_reference_file(file_url=fu, file_type=ft)

            parsed_results: list[tuple[dict, dict | None]] = []
            if parse_jobs:
                async with asyncio.TaskGroup() as tg:
                    tasks = [tg.create_task(_one_parse(rf, ft, fu)) for rf, ft, fu in parse_jobs]
                parsed_results = [t.result() for t in tasks]

            for idx, (rf, parsed) in enumerate(parsed_results, start=1):
                ft = rf.get("file_type") or "unknown"
                ins = (rf.get("instruction") or "").strip()
                blocks.append(f"### 参考资料{idx}（{ft}）")
                if ins:
                    blocks.append(f"使用说明：{ins}")
                if not parsed:
                    blocks.append("（未获取到解析结果，保留文件引用）")
                    continue

                summary = parsed.get("summary")
                if isinstance(summary, str) and summary.strip():
                    blocks.append(summary.strip())

                text_chunks = parsed.get("text_chunks") or []
                if isinstance(text_chunks, list):
                    for cidx, ch in enumerate(text_chunks[:6], start=1):
                        content = ch.get("content") if isinstance(ch, dict) else ""
                        if isinstance(content, str) and content.strip():
                            blocks.append(f"- 片段{cidx}: {content.strip()}")

                images = parsed.get("images") or []
                if isinstance(images, list):
                    for iidx, img in enumerate(images[:4], start=1):
                        if not isinstance(img, dict):
                            continue
                        img_url = (img.get("image_url") or "").strip()
                        if not img_url:
                            continue
                        ext = self._infer_extension(img_url, fallback=".png")
                        local_name = self._sanitize_filename(
                            f"ref_{idx}_img_{iidx}{ext}"
                        )
                        try:
                            await self._download_reference_image(
                                file_url=img_url, dest_path=image_dir / local_name
                            )
                            blocks.append(f"![参考资料{idx}-图片{iidx}]({local_name})")
                        except Exception:
                            continue

                key_frames = parsed.get("key_frames") or []
                if isinstance(key_frames, list):
                    for kidx, kf in enumerate(key_frames[:3], start=1):
                        if not isinstance(kf, dict):
                            continue
                        desc = (kf.get("description") or "").strip()
                        transcript = (kf.get("transcript") or "").strip()
                        if desc or transcript:
                            blocks.append(
                                f"- 视频关键帧{kidx}: {desc} {('；语音：' + transcript) if transcript else ''}".strip()
                            )

        return "\n\n".join([b for b in blocks if b.strip()])

    async def generate_deck(
        self,
        *,
        task_dir: Path,
        user_id: str,
        topic: str,
        description: str,
        total_pages: int,
        audience: str,
        global_style: str,
        session_id: str,
        teaching_elements: Optional[dict] = None,
        reference_files: Optional[list[dict]] = None,
    ) -> GeneratedDeck:
        """
        生成整个 PPT 初稿。
        其中 reference_files 目前至少支持 image：
        会下载到 `task_dir/images/` 并在 markdown 中以 `![...](filename)` 注入，
        让 PPTAgent 在生成时可以真正选择并嵌入这些图片。
        """
        self._load_template_once()
        assert self._template_presentation is not None
        assert self._template_slide_induction is not None

        # 由于 PPTGen.set_reference 会对 slide_induction 做 pop，避免污染缓存
        slide_induction = copy.deepcopy(self._template_slide_induction)

        ppt_agent = PPTAgent(
            language_model=self.language_model,
            vision_model=self.vision_model,
        )
        ppt_agent.set_reference(
            slide_induction=slide_induction,
            presentation=self._template_presentation,
            hide_small_pic_ratio=self.template_hidden_small_pic_ratio,
            keep_in_background=True,
        )

        # Document.from_markdown 会根据 markdown 中的 `![...](...)` 段落解析媒体。
        image_dir = task_dir / "images"
        image_dir.mkdir(parents=True, exist_ok=True)

        md = description_to_markdown(description, topic)
        external_md = await self._build_external_context_markdown(
            user_id=user_id,
            topic=topic,
            description=description,
            image_dir=image_dir,
            reference_files=reference_files,
        )
        if external_md:
            md += "\n\n" + external_md

        # Inject reference images into markdown so PPTAgent Document can parse them.
        if reference_files:
            image_refs: list[tuple[dict, str]] = []  # (rf, local_name)
            for i, rf in enumerate(reference_files, start=1):
                rf_type = (rf.get("file_type") or "").lower()
                if rf_type not in ("image", "images"):
                    continue
                file_url = (rf.get("file_url") or "").strip()
                if not file_url:
                    continue
                ext = self._infer_extension(file_url, fallback=".png")
                raw_local = rf.get("file_id") or f"img_{i:04d}"
                local_name = self._sanitize_filename(f"{raw_local}{ext}")
                image_refs.append((rf, local_name))

            if image_refs:
                md += "\n\n## 参考图片与素材\n"

                import asyncio

                async def _dl_one(rf: dict, local_name: str) -> None:
                    file_url = (rf.get("file_url") or "").strip()
                    if not file_url:
                        return
                    dest_path = image_dir / local_name
                    await self._download_reference_image(
                        file_url=file_url, dest_path=dest_path
                    )

                async with asyncio.TaskGroup() as tg:
                    for rf, local_name in image_refs:
                        tg.create_task(_dl_one(rf, local_name))

                for idx, (rf, local_name) in enumerate(image_refs, start=1):
                    instruction = (rf.get("instruction") or "").strip()
                    if instruction:
                        md += f"\n\n说明：{instruction}\n"
                    # IMPORTANT: image paragraph must be alone so regex `MARKDOWN_IMAGE_REGEX.match`
                    # can detect it as a Media block.
                    md += f"\n\n![参考图片{idx}]({local_name})\n"

        source_doc = await Document.from_markdown(
            markdown_content=md,
            language_model=self.language_model,
            vision_model=self.vision_model,
            image_dir=str(image_dir),
        )

        num_slides = None if total_pages <= 0 else int(total_pages)
        prs, _history = await ppt_agent.generate_pres(
            source_doc=source_doc,
            num_slides=num_slides,
            outline=None,
            image_dir=str(image_dir),
            dst_language=source_doc.language,
        )

        task_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = task_dir / "output.pptx"
        prs.save(str(pptx_path))

        slide_image_dir = task_dir / "renders"
        slide_image_dir.mkdir(parents=True, exist_ok=True)
        await ppt_to_images(str(pptx_path), str(slide_image_dir))

        slide_html: list[str] = [slide.to_html() for slide in prs.slides]

        return GeneratedDeck(
            task_dir=task_dir,
            pptx_path=pptx_path,
            slide_image_dir=slide_image_dir,
            slide_html=slide_html,
            slide_count=len(prs.slides),
        )

    async def edit_slide_html_llm(
        self,
        *,
        current_html: str,
        instruction: str,
        action_type: str,
        topic: str,
    ) -> str:
        """
        后续修改：仅用 AsyncLLM 改单页 HTML，不跑 PPTAgent.generate_pres。
        """
        prompt = (
            "你是课件幻灯片 HTML 编辑。只输出一段 HTML 片段"
            "（不要 DOCTYPE、不要 html/head/body 外壳），用于 16:9 幻灯片区域。\n"
            f"课程主题：{topic}\n"
            f"操作类型：{action_type}\n\n"
            f"【当前幻灯片 HTML】\n{current_html[:8000]}\n\n"
            f"【教师修改指令】\n{instruction}\n\n"
            "要求：教学向排版，可用内联 style；禁止 script、禁止 Markdown；"
            "只输出替换后的主体 HTML。"
        )
        raw = await self.language_model(
            prompt,
            system_message="只输出 HTML 片段，不要解释。",
            return_json=False,
        )
        if isinstance(raw, tuple):
            raw = raw[0]
        text = str(raw).strip()
        if text.startswith("```"):
            lines = text.split("\n")
            if lines and lines[0].startswith("```"):
                lines = lines[1:]
            if lines and lines[-1].strip().startswith("```"):
                lines = lines[:-1]
            text = "\n".join(lines).strip()
        return text if text else current_html

    def refresh_slide_assets(self, task: object, page_ids: list[str]) -> None:
        """
        将已更新的 py_code 同步到 renders/*.jpg，并尽力同步 output.pptx 首文本框（不重建整册）。
        """
        from ppt_agent_service.task_manager import PageState, TaskState, utc_ms

        if not isinstance(task, TaskState) or not task.output_pptx_path:
            return
        pptx_path = task.output_pptx_path
        task_dir = pptx_path.parent
        for pid in page_ids:
            p = task.pages.get(pid)
            if not isinstance(p, PageState):
                continue
            inner = extract_slide_html_from_py(p.py_code)
            self._write_slide_preview_jpeg(task_dir, p.slide_index, inner)
            self._sync_pptx_slide_plain_text(pptx_path, p.slide_index, inner)
        task.last_update = utc_ms()

    def _write_slide_preview_jpeg(
        self, task_dir: Path, slide_index: int, html_inner: str
    ) -> None:
        try:
            from html2image import Html2Image
            from PIL import Image
        except ImportError:
            return
        renders = task_dir / "renders"
        renders.mkdir(parents=True, exist_ok=True)
        png_name = f"slide_{slide_index:04d}_tmp.png"
        jpg_name = f"slide_{slide_index:04d}.jpg"
        full_html = (
            "<!DOCTYPE html><html><head><meta charset=\"utf-8\"/>"
            "<style>body{margin:0;padding:28px;box-sizing:border-box;width:1280px;height:720px;"
            "font-family:'Microsoft YaHei','Segoe UI',sans-serif;background:#fff;}"
            "</style></head><body>"
            f"{html_inner}</body></html>"
        )
        try:
            hti = Html2Image(output_path=str(renders), size=(1280, 720))
            hti.screenshot(html_str=full_html, save_as=png_name)
            png_path = renders / png_name
            jpg_path = renders / jpg_name
            if png_path.exists():
                im = Image.open(png_path).convert("RGB")
                im.save(jpg_path, "JPEG", quality=88)
                png_path.unlink(missing_ok=True)
        except Exception:
            pass

    def execute_render_job(
        self, task_dir: Path, slide_index: int, py_code: str
    ) -> tuple[bool, str]:
        """
        §3.8.3：从 py_code 提取 HTML 并写出 renders/slide_XXXX.jpg（需 html2image + PIL）。
        """
        try:
            from html2image import Html2Image
            from PIL import Image
        except ImportError as e:
            return False, f"html2image/PIL 未安装: {e}"
        try:
            inner = extract_slide_html_from_py(py_code)
        except Exception as e:
            return False, f"解析 py_code: {e}"
        renders = task_dir / "renders"
        renders.mkdir(parents=True, exist_ok=True)
        png_name = f"slide_{slide_index:04d}_tmp.png"
        jpg_name = f"slide_{slide_index:04d}.jpg"
        full_html = (
            "<!DOCTYPE html><html><head><meta charset=\"utf-8\"/>"
            "<style>body{margin:0;padding:28px;box-sizing:border-box;width:1280px;height:720px;"
            "font-family:'Microsoft YaHei','Segoe UI',sans-serif;background:#fff;}"
            "</style></head><body>"
            f"{inner}</body></html>"
        )
        png_path = renders / png_name
        jpg_path = renders / jpg_name
        try:
            hti = Html2Image(output_path=str(renders), size=(1280, 720))
            hti.screenshot(html_str=full_html, save_as=png_name)
            if png_path.exists():
                im = Image.open(png_path).convert("RGB")
                im.save(jpg_path, "JPEG", quality=88)
                png_path.unlink(missing_ok=True)
        except Exception as e:
            return False, str(e)
        if not jpg_path.exists():
            return False, "预览图未生成"
        return True, ""

    def _sync_pptx_slide_plain_text(
        self, pptx_path: Path, slide_index: int, html: str
    ) -> None:
        try:
            from pptx import Presentation
        except ImportError:
            return
        plain = re.sub(r"(?is)<script.*?>.*?</script>", " ", html)
        plain = re.sub(r"<[^>]+>", " ", plain)
        plain = " ".join(plain.split())
        if len(plain) > 4000:
            plain = plain[:4000] + "…"
        try:
            prs = Presentation(str(pptx_path))
            if slide_index < 1 or slide_index > len(prs.slides):
                return
            slide = prs.slides[slide_index - 1]
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = plain
                    break
            prs.save(str(pptx_path))
        except Exception:
            pass

